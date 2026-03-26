#!/usr/bin/env python3
"""HTML-to-PPTX converter.

Converts an HTML file (or a directory of HTML slide files) to a PPTX
presentation by rendering each page to a PNG image via Puppeteer and
embedding the images using python-pptx.

Node.js and Puppeteer are OPTIONAL — if unavailable, this module raises
ImportError with installation instructions.
"""

from __future__ import annotations

import argparse
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

from .renderer import render_html_to_image

# 16:9 slide dimensions (same constants as slide_factory.py)
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)


def convert(html_path: str, output_pptx: str) -> None:
    """Convert an HTML file to a PPTX presentation.

    Renders the HTML to a full-slide PNG screenshot via Puppeteer, then
    embeds the image into a new python-pptx presentation with 16:9 geometry.

    Parameters
    ----------
    html_path:
        Path to the HTML file to convert.  For multi-slide decks, pass the
        entry-point HTML (navigation or index page) — each distinct HTML file
        in the same directory will be rendered as a separate slide if
        ``html_path`` is a glob pattern (``*.html``).
    output_pptx:
        Destination ``.pptx`` file path.

    Raises
    ------
    ImportError
        If Node.js or Puppeteer is not found on PATH.
    FileNotFoundError
        If *html_path* does not exist.
    RuntimeError
        If Puppeteer rendering fails.
    """
    source = Path(html_path)

    # Collect HTML files to render
    if source.is_dir():
        html_files = sorted(source.glob("*.html"))
    elif "*" in str(source) or "?" in str(source):
        html_files = sorted(source.parent.glob(source.name))
    else:
        html_files = [source]

    if not html_files:
        raise FileNotFoundError(f"No HTML files found at: {html_path}")

    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H

    blank_layout = prs.slide_layouts[6]  # blank layout

    with tempfile.TemporaryDirectory() as tmpdir:
        for idx, hfile in enumerate(html_files):
            img_path = str(Path(tmpdir) / f"slide_{idx:04d}.png")
            render_html_to_image(str(hfile), img_path)

            slide = prs.slides.add_slide(blank_layout)
            _embed_image_fullslide(slide, img_path, prs)

    output_path = Path(output_pptx)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def _embed_image_fullslide(slide, image_path: str, prs: Presentation) -> None:
    """Embed *image_path* as a full-bleed image on *slide*."""
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    slide.shapes.add_picture(
        image_path,
        Emu(0),
        Emu(0),
        width=Emu(slide_w),
        height=Emu(slide_h),
    )


def main() -> int:
    parser = argparse.ArgumentParser(
        prog="html2pptx",
        description=__doc__,
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "html_path",
        help=("HTML file, directory, or glob pattern to convert.  Directories render every *.html file as a slide."),
    )
    parser.add_argument(
        "output_pptx",
        help="Output .pptx file path.",
    )
    args = parser.parse_args()

    try:
        convert(args.html_path, args.output_pptx)
        print(f"Saved: {args.output_pptx}")
        return 0
    except ImportError as exc:
        print(f"Dependency error: {exc}", file=sys.stderr)
        return 2
    except (FileNotFoundError, RuntimeError) as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    sys.exit(main())
