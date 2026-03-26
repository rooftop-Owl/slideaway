#!/usr/bin/env python3
"""Validate the structural integrity of a PPTX file.

Performs lightweight heuristic checks using python-pptx:

* The file can be opened (not corrupt).
* The presentation contains at least one slide.
* Slide dimensions are set (not None).
* No corrupt inline XML (detected via python-pptx parse errors).

This module provides a fast sanity check after the unpack → edit → pack
round-trip.  It does **not** enforce style conventions (aspect ratio, fonts,
footers, etc.) — those are handled by
:func:`modules.slides.tools.slide_factory.validate_pptx_structure`.

Example usage::

    from modules.slides.tools.pptx_editor import validate_pptx_xml

    result = validate_pptx_xml("deck.pptx")
    if result["valid"]:
        print("OK")
    else:
        for err in result["errors"]:
            print(f"  FAIL: {err}")

CLI::

    python3 validate.py deck.pptx
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path


def validate_pptx_xml(pptx_path: str) -> dict[str, bool | list[str]]:
    """Check structural integrity of a PPTX file.

    Args:
        pptx_path: Path to the ``.pptx`` file to validate.

    Returns:
        A dict with two keys:

        ``"valid"`` (bool)
            ``True`` if all checks pass.

        ``"errors"`` (list[str])
            Human-readable error messages; empty when *valid* is ``True``.

    Note:
        This function never raises — all exceptions are caught and returned as
        error strings so that callers can act on the result dict uniformly.

    Example::

        result = validate_pptx_xml("/tmp/repacked.pptx")
        assert result["valid"], result["errors"]
    """
    errors: list[str] = []
    path = Path(pptx_path)

    # ------------------------------------------------------------------ #
    # Check 1: file exists
    # ------------------------------------------------------------------ #
    if not path.exists():
        errors.append(f"File not found: {pptx_path}")
        return {"valid": False, "errors": errors}

    # ------------------------------------------------------------------ #
    # Check 2: importable — catch missing python-pptx gracefully
    # ------------------------------------------------------------------ #
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.exc import PackageNotFoundError  # type: ignore[import-untyped]
    except ImportError as exc:
        errors.append(f"python-pptx not installed: {exc}")
        return {"valid": False, "errors": errors}

    # ------------------------------------------------------------------ #
    # Check 3: file can be opened (not corrupt ZIP or corrupt OOXML)
    # ------------------------------------------------------------------ #
    try:
        prs = Presentation(str(path))
    except PackageNotFoundError as exc:
        errors.append(f"Cannot open PPTX (corrupt or not a PPTX): {exc}")
        return {"valid": False, "errors": errors}
    except Exception as exc:  # noqa: BLE001 — intentional broad catch
        errors.append(f"Unexpected error opening PPTX: {exc}")
        return {"valid": False, "errors": errors}

    # ------------------------------------------------------------------ #
    # Check 4: at least one slide
    # ------------------------------------------------------------------ #
    slide_count = len(prs.slides)
    if slide_count == 0:
        errors.append("Presentation has no slides")

    # ------------------------------------------------------------------ #
    # Check 5: slide dimensions are present
    # ------------------------------------------------------------------ #
    if prs.slide_width is None or prs.slide_height is None:
        errors.append("Slide dimensions are not set (slide_width or slide_height is None)")

    # ------------------------------------------------------------------ #
    # Check 6: no corrupt XML in slide shapes (iterate to trigger parse)
    # ------------------------------------------------------------------ #
    for idx, slide in enumerate(prs.slides, 1):
        try:
            # Accessing shapes triggers XML element parsing
            _ = list(slide.shapes)
        except Exception as exc:  # noqa: BLE001
            errors.append(f"Slide {idx}: corrupt XML detected — {exc}")

    valid = len(errors) == 0
    return {"valid": valid, "errors": errors}


def main() -> int:
    parser = argparse.ArgumentParser(
        prog="validate",
        description=(
            "Validate the structural integrity of a PPTX file.\n\n"
            "Checks: file exists, can be opened, has slides, slide dimensions\n"
            "are present, and slide XML is parseable."
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "Example:\n"
            "  python3 validate.py /tmp/repacked.pptx\n"
            "\n"
            "Exit codes:\n"
            "  0 — all checks passed\n"
            "  1 — one or more checks failed"
        ),
    )
    parser.add_argument("pptx_path", help="Path to the .pptx file to validate")
    args = parser.parse_args()

    result = validate_pptx_xml(args.pptx_path)
    valid: bool = result["valid"]  # type: ignore[assignment]
    errors: list[str] = result["errors"]  # type: ignore[assignment]

    if valid:
        print(f"PASS: '{args.pptx_path}' is structurally valid")
        return 0
    else:
        print(f"FAIL: '{args.pptx_path}' failed validation:")
        for err in errors:
            print(f"  ERROR: {err}")
        return 1


if __name__ == "__main__":
    sys.exit(main())
