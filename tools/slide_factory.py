#!/usr/bin/env python3

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Sequence, TypedDict
import re

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


SLIDE_W_16_9 = Inches(13.333)
SLIDE_H_16_9 = Inches(7.5)


@dataclass
class Theme:
    primary: RGBColor = field(default_factory=lambda: RGBColor(0x1A, 0x3A, 0x5C))
    accent: RGBColor = field(default_factory=lambda: RGBColor(0x2E, 0x86, 0xAB))
    surface: RGBColor = field(default_factory=lambda: RGBColor(0xF0, 0xF4, 0xF8))
    text: RGBColor = field(default_factory=lambda: RGBColor(0x33, 0x33, 0x33))
    text_muted: RGBColor = field(default_factory=lambda: RGBColor(0x66, 0x66, 0x66))
    subtitle_text: RGBColor = field(default_factory=lambda: RGBColor(0xCC, 0xDD, 0xEE))
    font_name: str = "Calibri"

    @classmethod
    def from_style(cls, name: str) -> "Theme":
        """Return a Theme preset by name. Raises ValueError for unknown names."""
        if name not in STYLE_PRESETS:
            available = ", ".join(sorted(STYLE_PRESETS.keys()))
            raise ValueError(f"Unknown style '{name}'. Available: {available}")
        return STYLE_PRESETS[name]

    @classmethod
    def list_styles(cls) -> list[str]:
        """Return sorted list of available style preset names."""
        return sorted(STYLE_PRESETS.keys())


STYLE_PRESETS: dict[str, Theme] = {
    "corporate": Theme(
        primary=RGBColor(0x1A, 0x3A, 0x5C),
        accent=RGBColor(0xC9, 0xA8, 0x4C),
        surface=RGBColor(0xFF, 0xFF, 0xFF),
        text=RGBColor(0x1A, 0x1A, 0x1A),
        text_muted=RGBColor(0x6B, 0x6B, 0x6B),
        subtitle_text=RGBColor(0x2E, 0x6D, 0xA4),
        font_name="Calibri",
    ),
    "academic": Theme(
        primary=RGBColor(0x1A, 0x3A, 0x5C),
        accent=RGBColor(0x8B, 0x69, 0x14),
        surface=RGBColor(0xFD, 0xFC, 0xF8),
        text=RGBColor(0x1A, 0x1A, 0x1A),
        text_muted=RGBColor(0x4A, 0x4A, 0x4A),
        subtitle_text=RGBColor(0x2D, 0x59, 0x86),
        font_name="Georgia",
    ),
    "creative": Theme(
        primary=RGBColor(0xFF, 0x3B, 0x5C),
        accent=RGBColor(0xFF, 0xD7, 0x00),
        surface=RGBColor(0xFF, 0xFF, 0xFF),
        text=RGBColor(0x1A, 0x1A, 0x1A),
        text_muted=RGBColor(0x55, 0x55, 0x55),
        subtitle_text=RGBColor(0x00, 0xE5, 0xFF),
        font_name="Calibri",
    ),
    "minimalist": Theme(
        primary=RGBColor(0x00, 0x71, 0xE3),
        accent=RGBColor(0x00, 0x71, 0xE3),
        surface=RGBColor(0xFF, 0xFF, 0xFF),
        text=RGBColor(0x1D, 0x1D, 0x1F),
        text_muted=RGBColor(0x6E, 0x6E, 0x73),
        subtitle_text=RGBColor(0x6E, 0x6E, 0x73),
        font_name="Calibri",
    ),
    "bold": Theme(
        primary=RGBColor(0xFF, 0x51, 0x00),
        accent=RGBColor(0xFF, 0xFF, 0xFF),
        surface=RGBColor(0x0A, 0x0A, 0x0A),
        text=RGBColor(0xFF, 0xFF, 0xFF),
        text_muted=RGBColor(0xFF, 0x51, 0x00),
        subtitle_text=RGBColor(0xFF, 0xFF, 0xFF),
        font_name="Calibri",
    ),
    "elegant": Theme(
        primary=RGBColor(0xC9, 0xA8, 0x4C),
        accent=RGBColor(0xE8, 0xD5, 0xA3),
        surface=RGBColor(0x0A, 0x0F, 0x1E),
        text=RGBColor(0xF2, 0xEC, 0xD6),
        text_muted=RGBColor(0xA0, 0x98, 0x80),
        subtitle_text=RGBColor(0xA0, 0x84, 0x5C),
        font_name="Calibri",
    ),
    "tech": Theme(
        primary=RGBColor(0x0E, 0xA5, 0xE9),
        accent=RGBColor(0xF5, 0x9E, 0x0B),
        surface=RGBColor(0x0F, 0x17, 0x2A),
        text=RGBColor(0xF1, 0xF5, 0xF9),
        text_muted=RGBColor(0x94, 0xA3, 0xB8),
        subtitle_text=RGBColor(0x38, 0xBD, 0xF8),
        font_name="Calibri",
    ),
    "warm": Theme(
        primary=RGBColor(0xC1, 0x7A, 0x40),
        accent=RGBColor(0xD4, 0xA0, 0x55),
        surface=RGBColor(0xF5, 0xED, 0xD9),
        text=RGBColor(0x2C, 0x18, 0x10),
        text_muted=RGBColor(0x7A, 0x55, 0x40),
        subtitle_text=RGBColor(0x8B, 0x5E, 0x3C),
        font_name="Calibri",
    ),
    "dark": Theme(
        primary=RGBColor(0x00, 0xFF, 0x88),
        accent=RGBColor(0x00, 0xD4, 0xFF),
        surface=RGBColor(0x05, 0x05, 0x10),
        text=RGBColor(0xE0, 0xE0, 0xF0),
        text_muted=RGBColor(0x88, 0x88, 0xBB),
        subtitle_text=RGBColor(0x7B, 0x00, 0xFF),
        font_name="Calibri",
    ),
    "nature": Theme(
        primary=RGBColor(0x3D, 0x52, 0x47),
        accent=RGBColor(0xC1, 0x7F, 0x5B),
        surface=RGBColor(0xF4, 0xF0, 0xEA),
        text=RGBColor(0x2A, 0x28, 0x25),
        text_muted=RGBColor(0x7A, 0x6E, 0x65),
        subtitle_text=RGBColor(0x6B, 0x8E, 0x7F),
        font_name="Calibri",
    ),
}

STYLE_PRESETS.update({
    "vibrant-pop": Theme(
        primary=RGBColor(0xFF, 0x6B, 0x35),
        accent=RGBColor(0x00, 0xC9, 0xA7),
        surface=RGBColor(0xFF, 0xFA, 0xF5),
        text=RGBColor(0x1A, 0x0A, 0x00),
        text_muted=RGBColor(0x7A, 0x4A, 0x30),
        subtitle_text=RGBColor(0x00, 0x9B, 0x7E),
        font_name="Calibri",
    ),
    "bold-contrast": Theme(
        primary=RGBColor(0x0D, 0x0D, 0x0D),
        accent=RGBColor(0xFF, 0x6D, 0x00),
        surface=RGBColor(0xF5, 0xF5, 0xF5),
        text=RGBColor(0x0D, 0x0D, 0x0D),
        text_muted=RGBColor(0x44, 0x44, 0x44),
        subtitle_text=RGBColor(0xFF, 0x6D, 0x00),
        font_name="Calibri",
    ),
    "electric-neon": Theme(
        primary=RGBColor(0x00, 0xE5, 0xFF),
        accent=RGBColor(0xFF, 0x00, 0xC8),
        surface=RGBColor(0x08, 0x08, 0x12),
        text=RGBColor(0xE8, 0xE8, 0xFF),
        text_muted=RGBColor(0x70, 0x70, 0xCC),
        subtitle_text=RGBColor(0x00, 0xE5, 0xFF),
        font_name="Calibri",
    ),
    "gradient-flow": Theme(
        primary=RGBColor(0xFF, 0x6E, 0xC7),
        accent=RGBColor(0x00, 0xD4, 0xFF),
        surface=RGBColor(0xFF, 0xB3, 0x47),
        text=RGBColor(0xFF, 0xFF, 0xFF),
        text_muted=RGBColor(0xCC, 0xCC, 0xFF),
        subtitle_text=RGBColor(0x7B, 0x61, 0xFF),
        font_name="Calibri",
    ),
    "dynamic-energy": Theme(
        primary=RGBColor(0xFF, 0x17, 0x44),
        accent=RGBColor(0xFF, 0xEA, 0x00),
        surface=RGBColor(0x0D, 0x0D, 0x0D),
        text=RGBColor(0xFF, 0xFF, 0xFF),
        text_muted=RGBColor(0xFF, 0xAB, 0x40),
        subtitle_text=RGBColor(0xFF, 0x6D, 0x00),
        font_name="Calibri",
    ),
    "academic-classic": Theme(
        primary=RGBColor(0x4A, 0x23, 0x5A),
        accent=RGBColor(0xB8, 0x86, 0x0B),
        surface=RGBColor(0xFD, 0xFC, 0xF8),
        text=RGBColor(0x1A, 0x1A, 0x1A),
        text_muted=RGBColor(0x4A, 0x4A, 0x4A),
        subtitle_text=RGBColor(0x6A, 0x3A, 0x80),
        font_name="Georgia",
    ),
    "research-paper": Theme(
        primary=RGBColor(0x1A, 0x44, 0x80),
        accent=RGBColor(0xC0, 0x00, 0x00),
        surface=RGBColor(0xFF, 0xFF, 0xFF),
        text=RGBColor(0x00, 0x00, 0x00),
        text_muted=RGBColor(0x59, 0x59, 0x59),
        subtitle_text=RGBColor(0x2E, 0x75, 0xB6),
        font_name="Calibri",
    ),
    "scientific-data": Theme(
        primary=RGBColor(0x00, 0x72, 0xB2),
        accent=RGBColor(0x00, 0x9E, 0x73),
        surface=RGBColor(0xFF, 0xFF, 0xFF),
        text=RGBColor(0x1A, 0x1A, 0x1A),
        text_muted=RGBColor(0x66, 0x66, 0x66),
        subtitle_text=RGBColor(0xE6, 0x9F, 0x00),
        font_name="Calibri",
    ),
    "scholarly-dark": Theme(
        primary=RGBColor(0xC9, 0xA8, 0x4C),
        accent=RGBColor(0x8A, 0x73, 0x40),
        surface=RGBColor(0x1A, 0x12, 0x08),
        text=RGBColor(0xD4, 0xBF, 0x9A),
        text_muted=RGBColor(0x9A, 0x8A, 0x70),
        subtitle_text=RGBColor(0xD4, 0xBF, 0x9A),
        font_name="Georgia",
    ),
    "publication-ready": Theme(
        primary=RGBColor(0x00, 0x33, 0x66),
        accent=RGBColor(0xCC, 0x33, 0x00),
        surface=RGBColor(0xFF, 0xFF, 0xFF),
        text=RGBColor(0x1A, 0x1A, 0x1A),
        text_muted=RGBColor(0x4D, 0x4D, 0x4D),
        subtitle_text=RGBColor(0x33, 0x66, 0x99),
        font_name="Calibri",
    ),
    "warm-earth": Theme(
        primary=RGBColor(0xA0, 0x4A, 0x2A),
        accent=RGBColor(0x6B, 0x8E, 0x5E),
        surface=RGBColor(0xF7, 0xF0, 0xE6),
        text=RGBColor(0x2A, 0x14, 0x08),
        text_muted=RGBColor(0x7A, 0x4A, 0x30),
        subtitle_text=RGBColor(0x5A, 0x7A, 0x4A),
        font_name="Calibri",
    ),
    "natural-wood": Theme(
        primary=RGBColor(0x6B, 0x42, 0x26),
        accent=RGBColor(0xC8, 0x90, 0x2E),
        surface=RGBColor(0xF0, 0xE6, 0xD3),
        text=RGBColor(0x2A, 0x1A, 0x0E),
        text_muted=RGBColor(0x7A, 0x55, 0x35),
        subtitle_text=RGBColor(0x8B, 0x63, 0x47),
        font_name="Calibri",
    ),
    "sage-garden": Theme(
        primary=RGBColor(0x62, 0x7B, 0x6A),
        accent=RGBColor(0xC1, 0x7C, 0x8A),
        surface=RGBColor(0xF2, 0xEF, 0xE9),
        text=RGBColor(0x2C, 0x33, 0x28),
        text_muted=RGBColor(0x7A, 0x80, 0x70),
        subtitle_text=RGBColor(0x8F, 0xAD, 0x96),
        font_name="Calibri",
    ),
    "sunset-glow": Theme(
        primary=RGBColor(0xFF, 0x70, 0x43),
        accent=RGBColor(0xF0, 0x62, 0x92),
        surface=RGBColor(0xFF, 0xF3, 0xE0),
        text=RGBColor(0x3E, 0x27, 0x23),
        text_muted=RGBColor(0x8D, 0x6E, 0x63),
        subtitle_text=RGBColor(0xFF, 0x98, 0x00),
        font_name="Calibri",
    ),
    "desert-sand": Theme(
        primary=RGBColor(0xC4, 0x81, 0x3A),
        accent=RGBColor(0x6B, 0x4C, 0x3B),
        surface=RGBColor(0xF7, 0xED, 0xD8),
        text=RGBColor(0x2E, 0x1B, 0x0E),
        text_muted=RGBColor(0x8C, 0x6B, 0x50),
        subtitle_text=RGBColor(0xA0, 0x78, 0x5A),
        font_name="Calibri",
    ),
    "midnight-premium": Theme(
        primary=RGBColor(0x0D, 0x1B, 0x3E),
        accent=RGBColor(0xD4, 0xAF, 0x37),
        surface=RGBColor(0x06, 0x0C, 0x1E),
        text=RGBColor(0xF0, 0xE8, 0xD0),
        text_muted=RGBColor(0x8A, 0x82, 0x6A),
        subtitle_text=RGBColor(0xD4, 0xAF, 0x37),
        font_name="Calibri",
    ),
    "dark-luxury": Theme(
        primary=RGBColor(0xB7, 0x6E, 0x79),
        accent=RGBColor(0xE8, 0xD5, 0xC4),
        surface=RGBColor(0x12, 0x12, 0x12),
        text=RGBColor(0xF5, 0xF0, 0xEB),
        text_muted=RGBColor(0x9A, 0x88, 0x78),
        subtitle_text=RGBColor(0x9C, 0x7B, 0x6E),
        font_name="Calibri",
    ),
    "obsidian": Theme(
        primary=RGBColor(0xFF, 0xFF, 0xFF),
        accent=RGBColor(0xFF, 0xFF, 0xFF),
        surface=RGBColor(0x00, 0x00, 0x00),
        text=RGBColor(0xFF, 0xFF, 0xFF),
        text_muted=RGBColor(0xB0, 0xB0, 0xB0),
        subtitle_text=RGBColor(0xE0, 0xE0, 0xE0),
        font_name="Calibri",
    ),
    "deep-ocean": Theme(
        primary=RGBColor(0x00, 0xB4, 0xD8),
        accent=RGBColor(0x90, 0xE0, 0xEF),
        surface=RGBColor(0x03, 0x04, 0x5E),
        text=RGBColor(0xCA, 0xF0, 0xF8),
        text_muted=RGBColor(0xAD, 0xE8, 0xF4),
        subtitle_text=RGBColor(0x00, 0x77, 0xB6),
        font_name="Calibri",
    ),
    "carbon-black": Theme(
        primary=RGBColor(0xE0, 0xE0, 0xE0),
        accent=RGBColor(0xFF, 0x6D, 0x00),
        surface=RGBColor(0x1A, 0x1A, 0x1A),
        text=RGBColor(0xEF, 0xEF, 0xEF),
        text_muted=RGBColor(0x9E, 0x9E, 0x9E),
        subtitle_text=RGBColor(0x9E, 0x9E, 0x9E),
        font_name="Calibri",
    ),
})
@dataclass
class Bullet:
    text: str
    level: int = 0
    bold_prefix: str | None = None


@dataclass
class ValidationResult:
    ok: bool
    info: list[str]
    warnings: list[str]
    errors: list[str]


class ChartSeries(TypedDict):
    name: str
    values: Sequence[float | int | None]


class ChartPayload(TypedDict):
    categories: Sequence[str]
    series: Sequence[ChartSeries]


class SlideFactory:
    def __init__(
        self,
        template_path: str | None = None,
        theme: Theme | None = None,
        force_widescreen: bool = False,
    ) -> None:
        self.prs = Presentation(template_path) if template_path else Presentation()
        if template_path is None or force_widescreen:
            self.prs.slide_width = SLIDE_W_16_9
            self.prs.slide_height = SLIDE_H_16_9
        self.theme = theme or Theme()
        self.margin = Inches(0.6)
        self.header_height = Inches(1.0)
        self.body_top = Inches(1.5)

    def _slide_width(self) -> int:
        width = self.prs.slide_width
        if width is None:
            raise ValueError("Presentation width is not set")
        return int(width)

    def _slide_height(self) -> int:
        height = self.prs.slide_height
        if height is None:
            raise ValueError("Presentation height is not set")
        return int(height)

    @property
    def body_width(self) -> Emu:
        return Emu(self._slide_width() - 2 * int(self.margin))

    @property
    def body_height(self) -> Emu:
        return Emu(self._slide_height() - int(self.body_top) - int(Inches(0.5)))

    def _add_bg(self, slide, color: RGBColor) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_header_bar(self, slide, title_text: str, subtitle: str | None = None) -> None:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(0),
            Emu(0),
            Emu(self._slide_width()),
            self.header_height,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme.primary
        bar.line.fill.background()

        title_box = slide.shapes.add_textbox(self.margin, Inches(0.15), self.body_width, Inches(0.7))
        title_tf = title_box.text_frame
        title_tf.clear()
        p = title_tf.paragraphs[0]
        r = p.add_run()
        r.text = title_text
        r.font.name = self.theme.font_name
        r.font.size = Pt(28)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        if subtitle:
            sub_box = slide.shapes.add_textbox(self.margin, Inches(0.6), self.body_width, Inches(0.35))
            sub_tf = sub_box.text_frame
            sub_tf.clear()
            p2 = sub_tf.paragraphs[0]
            r2 = p2.add_run()
            r2.text = subtitle
            r2.font.name = self.theme.font_name
            r2.font.size = Pt(15)
            r2.font.color.rgb = self.theme.subtitle_text

    def _add_footer(self, slide, slide_num: int, total_slides: int) -> None:
        footer = slide.shapes.add_textbox(
            Emu(self._slide_width() - int(Inches(1.2))),
            Emu(self._slide_height() - int(Inches(0.4))),
            Inches(1.0),
            Inches(0.3),
        )
        tf = footer.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = f"{slide_num} / {total_slides}"
        r.font.name = self.theme.font_name
        r.font.size = Pt(11)
        r.font.color.rgb = self.theme.text_muted

    def _add_bullets(self, slide, bullets: Sequence[Bullet]) -> None:
        box = slide.shapes.add_textbox(self.margin, self.body_top, self.body_width, self.body_height)
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True

        for i, item in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.level = item.level
            p.space_after = Pt(6)
            p.space_before = Pt(2)
            size = [Pt(20), Pt(16), Pt(14)][min(item.level, 2)]

            if item.bold_prefix:
                prefix = item.bold_prefix
                text = item.text
                rest = text[len(prefix) :] if text.startswith(prefix) else text

                r1 = p.add_run()
                r1.text = prefix
                r1.font.bold = True
                r1.font.name = self.theme.font_name
                r1.font.size = size
                r1.font.color.rgb = self.theme.text

                if rest:
                    r2 = p.add_run()
                    r2.text = rest
                    r2.font.name = self.theme.font_name
                    r2.font.size = size
                    r2.font.color.rgb = self.theme.text
            else:
                r = p.add_run()
                r.text = item.text
                r.font.name = self.theme.font_name
                r.font.size = size
                r.font.color.rgb = self.theme.text

    def add_title_slide(
        self,
        title: str,
        subtitle: str,
        author: str,
        date: str,
        region: str | None = None,
    ):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, self.theme.primary)

        title_box = slide.shapes.add_textbox(self.margin, Inches(1.8), self.body_width, Inches(1.2))
        title_tf = title_box.text_frame
        title_tf.clear()
        p = title_tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.name = self.theme.font_name
        r.font.size = Pt(54)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        subtitle_box = slide.shapes.add_textbox(self.margin, Inches(3.0), self.body_width, Inches(1.0))
        subtitle_tf = subtitle_box.text_frame
        subtitle_tf.clear()
        p2 = subtitle_tf.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = self.theme.font_name
        r2.font.size = Pt(28)
        r2.font.color.rgb = self.theme.subtitle_text

        if region:
            region_box = slide.shapes.add_textbox(self.margin, Inches(3.6), self.body_width, Inches(0.6))
            region_tf = region_box.text_frame
            region_tf.clear()
            p3 = region_tf.paragraphs[0]
            p3.alignment = PP_ALIGN.LEFT
            r3 = p3.add_run()
            r3.text = region
            r3.font.name = self.theme.font_name
            r3.font.size = Pt(24)
            r3.font.bold = True
            r3.font.color.rgb = self.theme.accent

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.margin, Inches(4.4), Inches(3.0), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme.accent
        line.line.fill.background()

        meta_box = slide.shapes.add_textbox(self.margin, Inches(4.8), self.body_width, Inches(0.8))
        meta_tf = meta_box.text_frame
        meta_tf.clear()
        p4 = meta_tf.paragraphs[0]
        r4 = p4.add_run()
        r4.text = author
        r4.font.name = self.theme.font_name
        r4.font.size = Pt(18)
        r4.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
        p5 = meta_tf.add_paragraph()
        r5 = p5.add_run()
        r5.text = date
        r5.font.name = self.theme.font_name
        r5.font.size = Pt(14)
        r5.font.color.rgb = RGBColor(0x88, 0x99, 0xAA)
        return slide

    def add_bullet_slide(
        self,
        title: str,
        bullets: Sequence[Bullet],
        subtitle: str | None = None,
        slide_num: int | None = None,
        total_slides: int | None = None,
    ):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_header_bar(slide, title, subtitle=subtitle)
        self._add_bullets(slide, bullets)
        if slide_num is not None and total_slides is not None:
            self._add_footer(slide, slide_num, total_slides)
        return slide

    def add_image_slide(
        self,
        title: str,
        image_path: str,
        caption: str | None = None,
        slide_num: int | None = None,
        total_slides: int | None = None,
    ):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_header_bar(slide, title)
        if slide_num is not None and total_slides is not None:
            self._add_footer(slide, slide_num, total_slides)

        max_w = self.body_width
        max_h = self.body_height

        try:
            from PIL import Image

            with Image.open(image_path) as img:
                iw, ih = img.size

            iw_emu = Emu(iw * 914400 // 96)
            ih_emu = Emu(ih * 914400 // 96)
            scale = min(int(max_w) / int(iw_emu), int(max_h) / int(ih_emu))
            pic_w = Emu(int(int(iw_emu) * scale))
            pic_h = Emu(int(int(ih_emu) * scale))
            pic_left = Emu((self._slide_width() - int(pic_w)) // 2)
            pic_top = Emu(int(self.body_top) + (int(max_h) - int(pic_h)) // 2)
            slide.shapes.add_picture(image_path, pic_left, pic_top, width=pic_w, height=pic_h)
        except Exception:
            pic = slide.shapes.add_picture(image_path, Emu(0), Emu(0), width=max_w)
            if int(pic.height) > int(max_h):
                ratio = int(max_h) / int(pic.height)
                pic.width = Emu(int(int(pic.width) * ratio))
                pic.height = max_h
            pic.left = Emu((self._slide_width() - int(pic.width)) // 2)
            pic.top = Emu(int(self.body_top) + (int(max_h) - int(pic.height)) // 2)

        if caption:
            cap_box = slide.shapes.add_textbox(
                self.margin,
                Emu(self._slide_height() - int(Inches(0.8))),
                self.body_width,
                Inches(0.25),
            )
            cap_tf = cap_box.text_frame
            cap_tf.clear()
            p = cap_tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = caption
            r.font.name = self.theme.font_name
            r.font.size = Pt(12)
            r.font.color.rgb = self.theme.text_muted
        return slide

    def add_chart_slide(
        self,
        title: str,
        chart_data: ChartPayload,
        slide_num: int | None = None,
        total_slides: int | None = None,
    ):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_header_bar(slide, title)
        if slide_num is not None and total_slides is not None:
            self._add_footer(slide, slide_num, total_slides)

        chart_obj = ChartData()
        chart_obj.categories = chart_data["categories"]
        for series in chart_data["series"]:
            chart_obj.add_series(series["name"], series["values"])

        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            self.margin,
            self.body_top,
            self.body_width,
            self.body_height,
            chart_obj,
        )
        return slide

    def save(self, path: str) -> None:
        output_path = Path(path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))

    @staticmethod
    def validate(path: str, allow_font_inheritance: bool = False) -> ValidationResult:
        return validate_pptx_structure(path, allow_font_inheritance=allow_font_inheritance)


def validate_pptx_structure(path: str, allow_font_inheritance: bool = False) -> ValidationResult:
    info: list[str] = []
    warnings: list[str] = []
    errors: list[str] = []

    prs = Presentation(path)
    slide_count = len(prs.slides)
    if slide_count == 0:
        errors.append("No slides in output")
        return ValidationResult(False, info, warnings, errors)

    width = prs.slide_width
    height = prs.slide_height
    if width is None or height is None:
        errors.append("Slide dimensions are missing")
        return ValidationResult(False, info, warnings, errors)

    slide_w = int(width)
    slide_h = int(height)
    ratio = slide_w / slide_h
    info.append(f"Slides: {slide_count}")
    if abs(ratio - (16 / 9)) > 0.05:
        errors.append(f"Aspect ratio is {ratio:.2f}, expected 16:9 (~1.78)")
    else:
        info.append(f"Aspect ratio: {ratio:.2f} (16:9)")

    first_text_blocks = 0
    for shape in prs.slides[0].shapes:
        text_frame = getattr(shape, "text_frame", None)
        if text_frame is None:
            continue
        if text_frame.text.strip():
            first_text_blocks += 1
    if first_text_blocks < 3:
        errors.append("Title slide must have separate title/subtitle/author-date text blocks")
    else:
        info.append(f"Title slide text blocks: {first_text_blocks}")

    footer_re = re.compile(r"^\s*(\d+)\s*/\s*(\d+)\s*$")
    missing_fonts: list[str] = []
    empty_slides: list[int] = []
    footer_lines: list[tuple[int, int, int]] = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        has_text = False
        has_non_text = False

        for shape in slide.shapes:
            left = int(getattr(shape, "left", 0) or 0)
            top = int(getattr(shape, "top", 0) or 0)
            shape_w = int(getattr(shape, "width", 0) or 0)
            shape_h = int(getattr(shape, "height", 0) or 0)
            if left < 0 or top < 0 or left + shape_w > slide_w or top + shape_h > slide_h:
                warnings.append(f"Slide {slide_idx}: shape appears outside slide bounds")

            shape_type = getattr(shape, "shape_type", None)
            if shape_type in {
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.CHART,
                MSO_SHAPE_TYPE.TABLE,
                MSO_SHAPE_TYPE.AUTO_SHAPE,
            }:
                has_non_text = True

            text_frame = getattr(shape, "text_frame", None)
            if text_frame is None:
                continue

            text = text_frame.text.strip()
            if text:
                has_text = True
                m = footer_re.match(text)
                if m:
                    footer_lines.append((slide_idx, int(m.group(1)), int(m.group(2))))

            for para in text_frame.paragraphs:
                for run in para.runs:
                    run_text = run.text or ""
                    if run_text.strip() and run.font.name is None:
                        missing_fonts.append(f"Slide {slide_idx}: '{run_text[:24]}'")

        if not has_text and not has_non_text:
            empty_slides.append(slide_idx)

    if empty_slides:
        errors.append(f"Empty slides detected: {empty_slides}")

    if missing_fonts:
        msg = f"Runs with font=None detected ({len(missing_fonts)}). Example: {missing_fonts[0]}"
        if allow_font_inheritance:
            warnings.append(msg + " (allowed by --allow-font-inheritance)")
        else:
            errors.append(msg)

    if footer_lines:
        by_slide = {slide_idx: (shown_num, shown_total) for slide_idx, shown_num, shown_total in footer_lines}
        if len(footer_lines) == slide_count - 1 and set(range(1, slide_count + 1)) - set(by_slide) == {1}:
            for slide_idx, (shown_num, shown_total) in by_slide.items():
                if shown_num != slide_idx or shown_total != slide_count:
                    warnings.append(
                        f"Footer mismatch on slide {slide_idx}: '{shown_num} / {shown_total}', expected '{slide_idx} / {slide_count}'"
                    )
        elif len(footer_lines) != slide_count:
            warnings.append(f"Footer numbering present on {len(footer_lines)}/{slide_count} slides")
        else:
            for slide_idx, shown_num, shown_total in footer_lines:
                if shown_num != slide_idx or shown_total != slide_count:
                    warnings.append(
                        f"Footer mismatch on slide {slide_idx}: '{shown_num} / {shown_total}', expected '{slide_idx} / {slide_count}'"
                    )

    return ValidationResult(len(errors) == 0, info, warnings, errors)


__all__ = [
    "Bullet",
    "SlideFactory",
    "STYLE_PRESETS",
    "Theme",
    "ValidationResult",
    "validate_pptx_structure",
]
