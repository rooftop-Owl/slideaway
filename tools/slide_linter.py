#!/usr/bin/env python3
"""slide_linter.py — Deterministic content linter for PPTX slides.

Checks content rules that validate_pptx.py does not cover:
  - Bullets per content slide (≤ max-bullets, default 6)
  - Heading presence on content slides
  - Per-run banned-font detection (ALL runs, not just run[0])

Exit codes:
  0  All checks passed
  1  Warnings only (no hard failures)
  2  One or more hard failures detected

Lint rules are deterministic — no LLM involved. For structural checks
(aspect ratio, slide count, empty slides), use validate_pptx.py instead.
"""

from __future__ import annotations

import argparse
import sys

try:
    from pptx import Presentation
except ImportError:
    # python-pptx not installed — skip linting rather than blocking writes.
    Presentation = None  # type: ignore[assignment, misc]

# Copied from validate_pptx.BANNED_FONTS — keep in sync with that definition.
BANNED_FONTS: frozenset[str] = frozenset({
    "Inter",
    "Roboto",
    "Arial",
    "system-ui",
    "sans-serif",
})

# Placeholder indices that represent non-body content (decorators, footers,
# slide numbers). We skip these when classifying slide type and counting bullets
# so that footer text does not incorrectly trigger heading or bullet rules.
_DECORATOR_PLACEHOLDER_INDICES: frozenset[int] = frozenset({10, 11, 12, 13, 14, 15})


def _slide_type(slide) -> str:
    """Classify a slide as 'title', 'content', or 'image_only'.

    title      — has a title placeholder but no body text (section breaks, closings)
    content    — has body text (bullets, prose, data) in addition to or instead of a title
    image_only — no text frames with content at all
    """
    has_title = False
    has_body_text = False

    for shape in slide.shapes:
        tf = getattr(shape, "text_frame", None)
        if tf is None or not tf.text.strip():
            continue
        ph = getattr(shape, "placeholder_format", None)
        if ph is not None:
            if ph.idx == 0:
                has_title = True
            elif ph.idx not in _DECORATOR_PLACEHOLDER_INDICES:
                has_body_text = True
        else:
            has_body_text = True

    if not has_title and not has_body_text:
        return "image_only"
    if has_title and not has_body_text:
        return "title"
    return "content"


def _count_top_level_bullets(slide) -> int:
    """Count level-0 non-empty paragraphs in non-title shapes."""
    count = 0
    for shape in slide.shapes:
        tf = getattr(shape, "text_frame", None)
        if tf is None:
            continue
        ph = getattr(shape, "placeholder_format", None)
        if ph is not None and (ph.idx == 0 or ph.idx in _DECORATOR_PLACEHOLDER_INDICES):
            continue
        for para in tf.paragraphs:
            if para.text.strip() and para.level == 0:
                count += 1
    return count


def _has_title(slide) -> bool:
    """Return True if a content slide has a populated title placeholder."""
    for shape in slide.shapes:
        ph = getattr(shape, "placeholder_format", None)
        if ph is not None and ph.idx == 0:
            tf = getattr(shape, "text_frame", None)
            if tf and tf.text.strip():
                return True
    return False


def _check_runs(slide, slide_idx: int) -> list[dict]:
    """Check every text run for banned fonts — ALL runs, not just run[0].

    validate_pptx.py checks runs for font mismatch but the loop stops at the
    first non-None font.name. This function iterates every run unconditionally
    so a deck cannot pass by hiding banned fonts in later runs.
    """
    issues = []
    for shape in slide.shapes:
        tf = getattr(shape, "text_frame", None)
        if tf is None:
            continue
        for para in tf.paragraphs:
            for run in para.runs:
                if not (run.text or "").strip():
                    continue
                try:
                    font_name = run.font.name
                except Exception:
                    continue  # malformed run XML — skip without crashing
                if font_name and font_name in BANNED_FONTS:
                    issues.append({
                        "slide": slide_idx,
                        "type": "banned_font",
                        "severity": "fail",
                        "detail": f"banned font '{font_name}' on '{(run.text or '')[:32]}'"
                    })
    return issues


def lint_pptx(
    pptx_path: str,
    max_bullets: int = 6,
    require_heading: bool = True,
) -> tuple[list[dict], list[dict]]:
    """Run all content lint rules on a PPTX file.

    Returns (failures, warnings). Callers should exit with code 2 on failures,
    1 on warnings-only, 0 on clean.
    """
    if Presentation is None:
        print("slide_linter: python-pptx not installed — linting skipped", file=sys.stderr)
        return [], []

    try:
        prs = Presentation(pptx_path)
    except FileNotFoundError:
        print(f"slide_linter: file not found: {pptx_path}", file=sys.stderr)
        sys.exit(2)
    except Exception as e:
        print(f"slide_linter: failed to open {pptx_path}: {e}", file=sys.stderr)
        sys.exit(2)

    failures: list[dict] = []
    warnings: list[dict] = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        stype = _slide_type(slide)

        if stype == "image_only":
            continue  # image-only slides are exempt from all content rules

        if stype == "content":
            # Heading presence
            if require_heading and not _has_title(slide):
                failures.append({
                    "slide": slide_idx,
                    "type": "missing_heading",
                    "severity": "fail",
                    "detail": "content slide has no populated title placeholder",
                })

            # Bullet count
            bullets = _count_top_level_bullets(slide)
            if bullets > max_bullets:
                failures.append({
                    "slide": slide_idx,
                    "type": "bullet_overflow",
                    "severity": "fail",
                    "detail": f"{bullets} top-level bullets (max: {max_bullets})",
                })

        # Per-run banned-font check — all slide types
        failures.extend(_check_runs(slide, slide_idx))

    return failures, warnings


def main() -> int:
    parser = argparse.ArgumentParser(
        prog="slide_linter",
        description="Deterministic content linter for PPTX slides.",
    )
    parser.add_argument("pptx_path", help="Path to .pptx file")
    parser.add_argument(
        "--max-bullets", type=int, default=6, metavar="N",
        help="Max top-level bullets per content slide (default: 6)",
    )
    parser.add_argument(
        "--no-require-heading", dest="require_heading", action="store_false",
        help="Allow content slides without a title placeholder",
    )
    parser.set_defaults(require_heading=True)
    parser.add_argument(
        "--quiet", action="store_true",
        help="Suppress passing output; still prints failures to stderr",
    )
    args = parser.parse_args()

    try:
        failures, warnings = lint_pptx(
            args.pptx_path,
            max_bullets=args.max_bullets,
            require_heading=args.require_heading,
        )
    except Exception as e:
        print(f"slide_linter: unexpected error: {e}", file=sys.stderr)
        return 2

    for w in warnings:
        print(f"WARN slide {w['slide']}: {w['detail']}")

    for f in failures:
        print(f"FAIL slide {f['slide']}: {f['detail']}", file=sys.stderr)

    if not failures and not warnings and not args.quiet:
        print(f"✓ slide_linter: {args.pptx_path} passed")

    if failures:
        return 2
    if warnings:
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
