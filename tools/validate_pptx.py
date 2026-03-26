#!/usr/bin/env python3

from __future__ import annotations

import argparse
import sys
from pathlib import Path

repo_root = Path(__file__).resolve().parents[3]
if str(repo_root) not in sys.path:
    sys.path.insert(0, str(repo_root))

from modules.slides.tools.slide_factory import Theme, validate_pptx_structure


def _channel_diff(c1, c2) -> int:
    """Maximum absolute difference per channel between two RGBColor values."""
    return max(abs(int(c1.r) - int(c2.r)), abs(int(c1.g) - int(c2.g)), abs(int(c1.b) - int(c2.b)))


def _near_palette(color, palette, tolerance: int = 30) -> bool:
    """Return True if color is within tolerance (per channel) of any palette entry."""
    return any(_channel_diff(color, p) <= tolerance for p in palette)


# Fonts that signal generic AI-generated output when used as the primary deck font.
# Detected in text runs with explicit font.name set to one of these values.
BANNED_FONTS: frozenset[str] = frozenset({
    "Inter",
    "Roboto",
    "Arial",
    "system-ui",
    "sans-serif",
})

def validate_style_compliance(pptx_path: str, theme: "Theme") -> tuple[list[str], list[str]]:
    """Validate font consistency and color palette adherence against a Theme preset.

    Checks every text run for explicit font names that differ from theme.font_name,
    and every shape fill for solid RGB colors far from the theme palette
    (primary, accent, surface) using a per-channel tolerance of 30.

    Returns (warnings, errors). All style violations are reported as warnings.
    """
    from pptx import Presentation

    warnings: list[str] = []
    errors: list[str] = []

    palette = [theme.primary, theme.accent, theme.surface]
    prs = Presentation(pptx_path)

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            # --- Font check ---
            # Only flag runs that set an explicit font name differing from the style font.
            # Runs with font.name=None inherit from the master (handled by --allow-font-inheritance).
            text_frame = getattr(shape, "text_frame", None)
            if text_frame is not None:
                for para in text_frame.paragraphs:
                    for run in para.runs:
                        run_text = run.text or ""
                        if run_text.strip() and run.font.name is not None:
                            if run.font.name != theme.font_name:
                                warnings.append(
                                    f"Slide {slide_idx}: font mismatch "
                                    f"'{run.font.name}' (expected '{theme.font_name}') "
                                    f"on text '{run_text[:24]}'"
                                )
                            if run.font.name in BANNED_FONTS:
                                warnings.append(
                                    f"Slide {slide_idx}: banned font '{run.font.name}' detected "
                                    f"(AI-slop heuristic) on text '{run_text[:24]}'"
                                )

            # --- Color check ---
            # Check solid fills with explicit RGB colors against the theme palette.
            # Non-solid fills, theme colors, and inherited fills are silently skipped.
            fill = getattr(shape, "fill", None)
            if fill is None:
                continue
            try:
                if fill.type is None:
                    continue
                fore = fill.fore_color
                if fore is None or fore.type is None:
                    continue
                color = fore.rgb  # raises if not an explicit RGB color
                if not _near_palette(color, palette):
                    warnings.append(
                        f"Slide {slide_idx}: fill color RGB({color.r},{color.g},{color.b}) "
                        f"is far from '{theme.font_name}' style palette"
                    )
            except (AttributeError, TypeError, ValueError):
                pass  # skip theme colors, gradients, patterns, etc.

    return warnings, errors


def main() -> int:
    parser = argparse.ArgumentParser(prog="validate_pptx")
    parser.add_argument("pptx_path", help="Path to .pptx file")
    parser.add_argument(
        "--allow-font-inheritance",
        action="store_true",
        help="Downgrade font=None findings to warnings (template-driven md2pptx mode)",
    )
    parser.add_argument(
        "--style",
        default=None,
        metavar="NAME",
        help=(
            f"Validate font/color consistency against a named style preset. Available: {', '.join(Theme.list_styles())}"
        ),
    )
    parser.add_argument(
        "--check-placeholders",
        action="store_true",
        help=(
            "Scan slide text for unfilled placeholder strings (e.g. 'xxxx', 'lorem', '[Title]', "
            "'Click to add'). Reports each match as FAIL with slide number and matched text."
        ),
    )
    args = parser.parse_args()

    result = validate_pptx_structure(
        args.pptx_path,
        allow_font_inheritance=args.allow_font_inheritance,
    )

    for item in result.info:
        print(f"INFO: {item}")
    for item in result.warnings:
        print(f"WARN: {item}")
    for item in result.errors:
        print(f"FAIL: {item}")

    # Style-aware validation — opt-in via --style, never changes baseline behavior
    style_ok = True
    if args.style is not None:
        try:
            theme = Theme.from_style(args.style)
        except ValueError as exc:
            print(f"FAIL: {exc}")
            return 1
        style_warnings, style_errors = validate_style_compliance(args.pptx_path, theme)
        for item in style_warnings:
            print(f"WARN: {item}")
        for item in style_errors:
            print(f"FAIL: {item}")
        style_ok = len(style_errors) == 0
    placeholder_ok = True
    if args.check_placeholders:
        from pptx import Presentation
        import re

        PLACEHOLDER_PATTERNS = [
            r"xxxx",
            r"lorem",
            r"ipsum",
            r"\[Title\]",
            r"\[Subtitle\]",
            r"Click to add",
            r"Insert text",
            r"placeholder",
        ]
        combined = re.compile("|".join(PLACEHOLDER_PATTERNS), re.IGNORECASE)

        prs = Presentation(args.pptx_path)
        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                tf = getattr(shape, "text_frame", None)
                if tf is None:
                    continue
                text = tf.text
                if not text:
                    continue
                for match in combined.finditer(text):
                    snippet = text[max(0, match.start() - 10) : match.end() + 10].replace("\n", " ")
                    print(f"FAIL: Slide {slide_idx}: placeholder text found: '...{snippet}...'")
                    placeholder_ok = False

    if result.ok and style_ok and placeholder_ok:
        print("PASS: structural inspection succeeded")
        return 0

    print("FAIL: structural inspection failed")
    return 1


if __name__ == "__main__":
    sys.exit(main())
